from flask import Flask, request, send_file
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import collections
import io

app = Flask(__name__)

# Color Logic
COLOR_MAP = {
    'Strong Offshore Opp': "#329B24",
    'Possible Offshore Opp': "#FFCA1A",
    'Offshore Opp': "#0000FF",
    'Not Offshorable': "#FF0000",
    'None': "#000071",
    'Dept': "#4D4D4D"
}

def hex_to_rgb(hex_color):
    hex_color = str(hex_color).lstrip('#')
    if hex_color in ["nan", "None", ""]: return (0, 0, 113)
    return tuple(int(hex_color[i:i + 2], 16) for i in (0, 2, 4))

def send_to_back(shape, slide):
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)

def generate_ppt_buffer(df):
    """Core logic remains EXACTLY the same, but saves to memory instead of disk"""
    nodes = {}
    edges = collections.defaultdict(list)

    ceo_row = df[df['Reports to'].isna() | (df['Reports to'] == "")].iloc[0]
    ceo_name = ceo_row['Name']
    nodes['CEO'] = {'title': 'CEO', 'subtitle': ceo_name, 'color': COLOR_MAP['None'], 'is_count': False}

    managers = set(df['Reports to'].dropna().unique())

    depts = [str(d) for d in df['Department'].unique() if str(d) != 'nan' and str(d).strip() != '']
    for d in depts:
        dept_id = f"DEPT_{d}"
        nodes[dept_id] = {'title': d, 'subtitle': 'Department', 'color': COLOR_MAP['Dept'], 'is_count': False}
        edges['CEO'].append(dept_id)

    for _, row in df.iterrows():
        name = row['Name']
        if name == ceo_name: continue
        if name in managers:
            nodes[name] = {'title': row['Title'], 'subtitle': name,
                           'color': COLOR_MAP.get(row['Label'], COLOR_MAP['None']), 'is_count': False}
            manager = row['Reports to']
            dept = str(row['Department'])
            if manager == ceo_name and dept in depts:
                edges[f"DEPT_{dept}"].append(name)
            else:
                edges[manager].append(name)

    leaves_df = df[~df['Name'].isin(managers) & (df['Name'] != ceo_name)]
    grouped = leaves_df.groupby(['Reports to', 'Department', 'Title', 'Label'], dropna=False).size().reset_index(name='Count')

    for idx, row in grouped.iterrows():
        node_id = f"LEAF_{idx}"
        nodes[node_id] = {'title': row['Title'], 'subtitle': str(row['Count']),
                          'color': COLOR_MAP.get(row['Label'], COLOR_MAP['None']), 'is_count': True}
        manager = row['Reports to']
        dept = str(row['Department'])
        if manager == ceo_name and dept in depts:
            edges[f"DEPT_{dept}"].append(node_id)
        elif manager in nodes:
            edges[manager].append(node_id)

    # LAYOUT ALGORITHM
    base_box_w, base_box_h = Inches(1.6), Inches(0.65)
    gap_x = Inches(0.4) 
    gap_y = Inches(1.3) 

    leaf_widths = {}

    def calc_width(n):
        if not edges[n]:
            leaf_widths[n] = base_box_w
            return base_box_w
        total = sum(calc_width(c) for c in edges[n]) + (gap_x * (len(edges[n]) - 1))
        leaf_widths[n] = max(base_box_w, total)
        return leaf_widths[n]

    calc_width('CEO')
    coords = {}

    def assign_coords(n, x_start, y):
        if not edges[n]:
            coords[n] = (x_start + base_box_w / 2, y)
            return

        current_x = x_start
        child_centers = []
        for c in edges[n]:
            assign_coords(c, current_x, y + base_box_h + gap_y)
            child_centers.append(coords[c][0])
            current_x += leaf_widths[c] + gap_x

        if len(child_centers) == 1:
            parent_x = child_centers[0]
        else:
            parent_x = (child_centers[0] + child_centers[-1]) / 2

        coords[n] = (parent_x, y)

    assign_coords('CEO', Inches(0.5), Inches(0.5))

    max_x = max(x for x, y in coords.values()) + base_box_w
    max_y = max(y for x, y in coords.values()) + base_box_h + Inches(1)

    MAX_PPT = Inches(54)
    scale = min(MAX_PPT / max_x, MAX_PPT / max_y) if max_x > MAX_PPT or max_y > MAX_PPT else 1.0

    prs = Presentation()
    prs.slide_width = max(Inches(13.33), int(max_x * scale))
    prs.slide_height = max(Inches(7.5), int(max_y * scale))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    drawn_shapes = {}

    for n_id, (raw_x, raw_y) in coords.items():
        data = nodes[n_id]
        x, y = raw_x * scale, raw_y * scale
        w, h = base_box_w * scale, base_box_h * scale
        title_font = max(Pt(4), int(Pt(9) * scale))
        sub_font = max(Pt(3), int(Pt(8) * scale))

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - (w / 2), y, w, h)
        rect.fill.solid()
        r, g, b = hex_to_rgb(data['color'])
        rect.fill.fore_color.rgb = RGBColor(r, g, b)
        rect.line.color.rgb = RGBColor(255, 255, 255)

        p = rect.text_frame.paragraphs[0]
        p.text = str(data['title'])
        p.font.size, p.font.bold, p.font.color.rgb = title_font, True, RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

        is_count = data['is_count']
        sub_w, sub_h = (Inches(0.4) * scale, Inches(0.25) * scale) if is_count else (w * 0.9, Inches(0.25) * scale)
        sub_x = (x + (w / 2) - sub_w) if is_count else (x - (sub_w / 2))
        sub_y = y + h - (sub_h / 2)

        sub_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, sub_x, sub_y, sub_w, sub_h)
        sub_rect.fill.solid()
        sub_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)

        sp = sub_rect.text_frame.paragraphs[0]
        sp.text = str(data['subtitle'])
        sp.font.size, sp.font.color.rgb = sub_font, RGBColor(0, 0, 0)
        sp.alignment = PP_ALIGN.CENTER

        drawn_shapes[n_id] = rect

    for parent, children in edges.items():
        if parent not in drawn_shapes: continue
        p_shape = drawn_shapes[parent]
        for child in children:
            if child not in drawn_shapes: continue
            c_shape = drawn_shapes[child]

            conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, 0, 0, 0, 0)
            conn.begin_connect(p_shape, 2)
            conn.end_connect(c_shape, 0)
            conn.line.color.rgb = RGBColor(120, 120, 120)
            conn.line.width = Pt(1.5 * scale)
            send_to_back(conn, slide)

    # Magic happens here: Save to memory instead of disk
    memory_file = io.BytesIO()
    prs.save(memory_file)
    memory_file.seek(0)
    return memory_file

@app.route('/', methods=['GET'])
def health_check():
    return "Org Chart API is running! Send a POST request to /generate with an Excel file.", 200

@app.route('/generate', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return {"error": "No file part in the request"}, 400
        
    file = request.files['file']
    
    if file.filename == '':
        return {"error": "No selected file"}, 400

    try:
        # Read Excel directly from the uploaded file stream
        df = pd.read_excel(file)
        
        # Generate the PPTX buffer
        ppt_buffer = generate_ppt_buffer(df)
        
        # Return as a downloadable file
        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name="Perfect_Structure_OrgChart.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        return {"error": str(e)}, 500

# Vercel needs this
if __name__ == "__main__":
    app.run(debug=True)
